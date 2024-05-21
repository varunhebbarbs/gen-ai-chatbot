import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain_community.llms import OpenAI
import pandas as pd
from pptx import Presentation
import os
from docx import Document
import time as t


class SessionState:
    def __init__(self, **kwargs):
        self.__dict__.update(kwargs)
 
def process_documents(documents_folder):
    documents = []
    #Extract the text
    for file_name in os.listdir(documents_folder):
        file_path = os.path.join(documents_folder, file_name)
        if file_path.endswith(".pdf"):
            pdf_reader = PdfReader(file_path)
            for page in pdf_reader.pages:
                documents.append(page.extract_text())
        elif file_path.endswith(".txt"):
            with open(file_path, "r") as file:
                documents.append(file.read())
        elif file_path.endswith(".docx"):
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                documents.append(paragraph.text)
        elif file_path.endswith(".xlsx"):
            xls = pd.ExcelFile(file_path)
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                for col in df.columns:
                    documents.extend(df[col].dropna().astype(str).tolist())
        elif file_path.endswith(".pptx"):
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        documents.append(shape.text)
 
    text = "\n".join(documents)
    if not text:
        return None
    
    #Break it into chuncks
    char_text_splitter = CharacterTextSplitter(separator="\n", chunk_size=1000, chunk_overlap=200, length_function=len)
    text_chunks = char_text_splitter.split_text(text)
    if not text_chunks:
        return None
 
    embeddings = OpenAIEmbeddings()
    #FAST Accurate nearest neighbour library
    # Creating vector store, -embeddings (OpenAI), -initializing FAISS, -Store chunks $ embeddings
    docsearch = FAISS.from_texts(text_chunks, embeddings)
    return docsearch
 
def answer_query(query, docsearch, chain):
        docs = docsearch.similarity_search(query)
        #chain -> take the question, get relevant document, pass it to the LLM, generate the output
        response = chain.run(input_documents=docs, question=query)
        return response
    
 
MAX_FILE_SIZE_MB = 2
 
def main():
    st.header("DOCUMENT SEARCH BOT ")
    st.info("The Document Search Bot is an intelligent application designed to assist users in efficiently searching and managing documents.")
    
    session_state = SessionState(documents_folder=None, docsearch=None, response="", is_admin=False)
 
    # Authentication mechanism
    user_type = st.sidebar.radio("Select user type:", ["User","Admin"])
    if user_type == "Admin":
        # Admin password authentication
        admin_password = st.sidebar.text_input("Enter admin password:", type="password")
        password_entered = False

        if admin_password:
         password_entered = True

        if password_entered:
           if admin_password == "admin":
            session_state.is_admin = True
            st.sidebar.success("Authentication successful")
           else:
            st.sidebar.error("Incorrect password. Please try again!.")
            
    if session_state.is_admin:
        # Admin functionality
        uploaded_files = st.file_uploader("Upload your documents📁", accept_multiple_files=True, type=["pdf", "txt", "docx", "xlsx", "pptx"])
        documents_folder = "admin_uploaded_docs"
        if not os.path.exists(documents_folder):
            os.makedirs(documents_folder)
        for file in uploaded_files:
            # Check file size
            if len(file.getvalue()) > MAX_FILE_SIZE_MB * 1024 * 1024:
                st.error(f"File size exceeds {MAX_FILE_SIZE_MB} MB limit: {file.name}")
                continue
 
            with open(os.path.join(documents_folder, file.name), "wb") as f:
                f.write(file.getvalue())
 
        session_state.documents_folder = documents_folder
        session_state.docsearch = process_documents(documents_folder)
        with st.spinner("Wait for a moment!"):
            t.sleep(2)
        chain = load_qa_chain(OpenAI(), chain_type="stuff")
 
        # Display uploaded documents
        uploaded_documents = os.listdir(documents_folder)
        if uploaded_documents:
            # AI assistant chat interface
            st.subheader("AI Assistant🤖")
            st.caption("Ask your question below based on uploaded documents")
            # chat_history = st.subheader.empty()
            user_input = st.text_input("You:", key="user_input")
            if st.button("Send"):
              response = answer_query(user_input, session_state.docsearch, chain)
              with st.spinner("Wait for a moment!"):
                t.sleep(2)
              st.text("Bot:")
              st.write(response)
            
            # display and append checked files to "selected_documents" 
            st.sidebar.header("Uploaded Documents:")
            selected_documents = []
            for file_name in os.listdir(documents_folder):
                selected = st.sidebar.checkbox(file_name)
                if selected:
                    selected_documents.append(file_name)
 
            # Delete selected documents
            if st.sidebar.button("Delete Selected Documents"):
                for file_name in selected_documents:
                    file_path = os.path.join(documents_folder, file_name)
                    os.remove(file_path)
                st.sidebar.write("Documents deleted successfully. Click below to refresh the page")
        else:
            st.warning("Documents not uploaded yet.")

            
    else:
        # Normal user functionality
        documents_folder = "admin_uploaded_docs"  # Use the admin's uploaded documents
        if not os.path.exists(documents_folder):
            st.warning("No documents uploaded yet!. Please wait for the admin to upload documents.")
            return
            
 
        session_state.documents_folder = documents_folder
        session_state.docsearch = process_documents(documents_folder)
        chain = load_qa_chain(OpenAI(), chain_type="stuff")

        # Fetch documents list
        uploaded_documents = os.listdir(documents_folder)
        if uploaded_documents:
            # Display uploaded documents
            st.sidebar.header("Avaliable Documents for Search:")
            for file_name in os.listdir(documents_folder): 
                st.sidebar.write(file_name)

                # AI assistant chat interface
                st.subheader("AI Assistant🤖")
                st.caption("Ask your question below based on uploaded documents")
                # chat_history = st.subheader.empty()
                user_input = st.text_input("You:", key="user_input")
                if st.button("Send"):
                    response = answer_query(user_input, session_state.docsearch, chain)
                    with st.spinner("Wait for a moment!"):
                        t.sleep(2)
                    st.text("Bot:")
                    st.write(response)

        else:# If the folder is empty
            st.warning("Wait for the admin to upload documents")
            
    if st.sidebar.button("Refresh Page"):
            st.experimental_rerun()
 
if __name__ == '__main__':
    main()
